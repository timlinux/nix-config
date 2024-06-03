from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "NixOS Caching Mechanisms with Flakes for Docker Images"
subtitle.text = "Efficient and Reproducible Builds\nPresented by: [Your Name]"

# Slide 2: Introduction
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Introduction to NixOS and Flakes"
content.text = (
    "Brief overview of NixOS\n"
    "Introduction to flakes and their purpose\n"
    "Importance of caching in build processes"
)

# Slide 3: The Nix Store
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "The Nix Store"
content.text = (
    "Explanation of the Nix store (/nix/store)\n"
    "Immutable storage of build artifacts\n"
    "Path naming convention using hashes"
)

# Slide 4: Derivation Hashing
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Derivation Hashing"
content.text = (
    "What is a derivation in Nix?\n"
    "How derivations are hashed\n"
    "Impact of changes on derivation hashes"
)

# Slide 5: Local Caching Mechanism
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Local Caching Mechanism"
content.text = (
    "Checking for existing store paths\n"
    "Reusing cached artifacts\n"
    "Efficiency gains from local caching"
)

# Slide 6: Binary Caches
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Binary Caches"
content.text = (
    "What are binary caches?\n"
    "Querying remote caches for artifacts\n"
    "Example: cache.nixos.org\n"
    "Benefits of using binary caches"
)

# Slide 7: Incremental Builds
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Incremental Builds"
content.text = (
    "Concept of incremental builds\n"
    "Reusing unchanged parts of the build\n"
    "Time and resource savings"
)

# Slide 8: Flake Lock Files
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Flake Lock Files"
content.text = (
    "Purpose of the flake.lock file\n"
    "Ensuring reproducibility and consistency\n"
    "How lock files affect caching"
)

# Slide 9: Putting It All Together
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Putting It All Together"
content.text = (
    "Step-by-step example of building a Docker image\n"
    "How NixOS uses caching at each step\n"
    "Flowchart of the caching process"
)

# Slide 10: Benefits and Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Benefits and Conclusion"
content.text = (
    "Summary of caching benefits (speed, efficiency, consistency)\n"
    "Reproducibility in software development\n"
    "Final thoughts and Q&A"
)

# Slide 11: Questions and Answers
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Questions and Answers"
content.text = (
    "Open floor for audience questions\n"
    "Clarifications and additional explanations"
)

# Save the presentation
prs.save("/mnt/data/nixos_caching_flakes_docker.pptx")

import shutil
shutil.move("/mnt/data/nixos_caching_flakes_docker.pptx", "nixos_caching_flakes_docker.pptx")
